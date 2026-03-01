# ============================================================
# CONVERTER — verbatim in das generierte Script einfuegen
#
# Generiertes Script-Geruest:
#   TOPIC = "Thema aus $ARGUMENTS"
#   HTML_CONTENT = """..."""
#   [dieser Code — verbatim]
# ============================================================

import subprocess, sys, os, base64
from pathlib import Path

# === DEPENDENCIES ===
for pkg in ['python-pptx', 'beautifulsoup4']:
    subprocess.run([sys.executable, '-m', 'pip', 'install', pkg, '-q'], check=True)

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === CI-FARBEN ===
ORANGE = RGBColor(0xEA, 0x5B, 0x0B)
GOLD   = RGBColor(0xFD, 0xC3, 0x00)
DARK   = RGBColor(0x56, 0x56, 0x56)
MGRAY  = RGBColor(0x99, 0x99, 0x99)
BGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_MAP = {'orange': ORANGE, 'gold': GOLD, 'dark': DARK}

TPL_DIR    = Path("/Users/robinwenzel/AI Playground/Claude Code/claude-code-skills/skills/slide/templates")
LOGO_PATH  = TPL_DIR / "logo_mindsquare.png"
LOGO_WHITE = TPL_DIR / "logo_mindsquare_white.png"
GRAD_PATH  = TPL_DIR / "gradient_line.png"

# === HEADER-ELEMENTE ===
def add_logo(slide, dark_bg=False):
    lp = str(LOGO_WHITE if dark_bg else LOGO_PATH)
    slide.shapes.add_picture(lp, Inches(11.0), Inches(0.15), height=Inches(0.5))

def add_gradient_line(slide):
    slide.shapes.add_picture(str(GRAD_PATH), Inches(0.2), Inches(0.84), Inches(12.87), Pt(3))

def add_footer(slide, page_num):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(11); p.font.color.rgb = MGRAY
    p.font.name = 'Calibri'; p.alignment = PP_ALIGN.RIGHT

def add_slide_title(slide, title_el):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10.0), Inches(0.42))
    tf = tb.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]
    for child in title_el.children:
        run = p.add_run()
        if isinstance(child, NavigableString):
            t = str(child)
            if not t.strip(): continue
            run.text = t; run.font.color.rgb = DARK
        elif isinstance(child, Tag):
            run.text = child.get_text()
            run.font.color.rgb = ORANGE if 'accent' in child.get('class', []) else DARK
        else:
            continue
        run.font.name = 'Poppins'; run.font.bold = True; run.font.size = Pt(20)

# === ICON-HILFSFUNKTIONEN ===
def draw_icon_target(slide, cx, cy, size, color):
    for i, ratio in enumerate([1.0, 0.65, 0.3]):
        r = int(size * ratio)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - r/2), int(cy - r/2), r, r)
        sh.fill.solid(); sh.fill.fore_color.rgb = color if i % 2 == 0 else WHITE
        sh.line.fill.background()

def draw_icon_person(slide, cx, cy, size, color):
    hr = int(size * 0.35)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - hr/2), int(cy - size*0.5), hr, hr)
    head.fill.solid(); head.fill.fore_color.rgb = color; head.line.fill.background()
    bw, bh = int(size * 0.6), int(size * 0.45)
    body = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, int(cx - bw/2), int(cy - size*0.1), bw, bh)
    body.fill.solid(); body.fill.fore_color.rgb = color; body.line.fill.background()

def draw_icon_people(slide, cx, cy, size, color):
    for ox, c in zip([-size*0.35, 0, size*0.35], [MGRAY, color, MGRAY]):
        draw_icon_person(slide, int(cx + ox), cy, int(size * 0.75), c)

def draw_icon_chart(slide, cx, cy, size, color):
    bw  = int(size * 0.2); gap = int(size * 0.1)
    total_w = 3 * bw + 2 * gap; sx = int(cx - total_w / 2)
    for i, (hr, c) in enumerate(zip([0.4, 0.7, 1.0], [MGRAY, GOLD, color])):
        bh = int(size * hr); bx = sx + i * (bw + gap); by = int(cy + size*0.5 - bh)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
        bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background()

def draw_icon_simple(slide, cx, cy, size, color, shape_type):
    sh = slide.shapes.add_shape(shape_type, int(cx-size/2), int(cy-size/2), int(size), int(size))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

ICON_BUILDERS = {
    'target':    lambda s,cx,cy,sz,c: draw_icon_target(s,cx,cy,sz,c),
    'gear':      lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.GEAR_6),
    'lightning': lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.LIGHTNING_BOLT),
    'person':    lambda s,cx,cy,sz,c: draw_icon_person(s,cx,cy,sz,c),
    'people':    lambda s,cx,cy,sz,c: draw_icon_people(s,cx,cy,sz,c),
    'chart':     lambda s,cx,cy,sz,c: draw_icon_chart(s,cx,cy,sz,c),
    'warning':   lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.DIAMOND),
    'document':  lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.FLOWCHART_DOCUMENT),
    'compass':   lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.PENTAGON),
    'cycle':     lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.CIRCULAR_ARROW),
    'funnel':    lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.FUNNEL),
    'cube':      lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.CUBE),
    'check':     lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.FLOWCHART_PROCESS),
    'euro':      lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.OVAL),
    'calendar':  lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.FLOWCHART_DOCUMENT),
    'arrow':     lambda s,cx,cy,sz,c: draw_icon_simple(s,cx,cy,sz,c, MSO_SHAPE.RIGHT_ARROW),
}

# === DYNAMISCHE SCHRIFTGROESSEN ===
def calc_bullet_pt(avail_h, n_items, max_pt=15.0, min_pt=10.0):
    """Optimale Bullet-Schriftgroesse: passt n_items in avail_h (EMU) ein."""
    if n_items == 0:
        return max_pt
    avail_in = avail_h / 914400
    if n_items * (max_pt * 1.3 + 3) / 72 <= avail_in:
        return max_pt
    size = (avail_in * 72 / n_items - 3) / 1.3
    return max(min_pt, min(max_pt, size))

def calc_heading_pt(col_w, max_pt=15.5, min_pt=12.0):
    """Ueberschrift-Schriftgroesse basierend auf Spaltenbreite (EMU)."""
    w_in = col_w / 914400
    size = min_pt + (max_pt - min_pt) * max(0.0, (w_in - 2.8)) / 3.5
    return max(min_pt, min(max_pt, size))

def calc_kpi_value_pt(col_w, max_pt=42.0, min_pt=28.0):
    """KPI-Wert-Schriftgroesse basierend auf Spaltenbreite (EMU)."""
    w_in = col_w / 914400
    size = min_pt + (max_pt - min_pt) * max(0.0, (w_in - 2.5)) / 4.5
    return max(min_pt, min(max_pt, size))

# === RENDER-FUNKTIONEN ===
def render_card(slide, card_el, left, top, width, height, accent='orange'):
    accent_color = ACCENT_MAP.get(accent, ORANGE)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(width), int(height))
    card.fill.solid(); card.fill.fore_color.rgb = BGRAY
    card.line.color.rgb = accent_color; card.line.width = Pt(2)

    cursor_y = top + Inches(0.18)

    icon_el = card_el.select_one('[data-icon]')
    if icon_el:
        icon_name = icon_el.get('data-icon', '')
        icon_size = Inches(0.45)
        builder   = ICON_BUILDERS.get(icon_name)
        if builder:
            builder(slide, int(left + width/2), int(cursor_y + icon_size/2), int(icon_size), accent_color)
        cursor_y += icon_size + Inches(0.1)

    h3 = card_el.select_one('h3')
    if h3:
        tb = slide.shapes.add_textbox(
            int(left + Inches(0.15)), int(cursor_y),
            int(width - Inches(0.3)), int(Inches(0.4)))
        tf = tb.text_frame
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h3.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(calc_heading_pt(width)); run.font.color.rgb = DARK
        cursor_y += Inches(0.44)

    bullets = card_el.select('ul li')
    if bullets:
        avail_h  = top + height - cursor_y - Inches(0.15)
        bullet_pt = calc_bullet_pt(avail_h, len(bullets))
        space_pt  = max(1.0, bullet_pt * 0.18)
        tb = slide.shapes.add_textbox(
            int(left + Inches(0.2)), int(cursor_y),
            int(width - Inches(0.35)), int(avail_h))
        tf = tb.text_frame; tf.word_wrap = True
        for j, li in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run(); run.text = '\u2022  ' + li.get_text(strip=True)
            run.font.name = 'Calibri'; run.font.size = Pt(bullet_pt)
            run.font.color.rgb = DARK
            p.space_after = Pt(space_pt)

def render_kpi_card(slide, card_el, left, top, width, height):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(width), int(height))
    card.fill.solid(); card.fill.fore_color.rgb = BGRAY
    card.line.color.rgb = ORANGE; card.line.width = Pt(2)

    cursor_y = top + Inches(0.2)

    icon_el = card_el.select_one('[data-icon]')
    if icon_el:
        icon_size = Inches(0.4)
        builder   = ICON_BUILDERS.get(icon_el.get('data-icon', ''))
        if builder:
            builder(slide, int(left + width/2), int(cursor_y + icon_size/2), int(icon_size), ORANGE)
        cursor_y += icon_size + Inches(0.1)

    value_el = card_el.select_one('.kpi-value')
    if value_el:
        ac  = value_el.get('data-accent', 'orange')
        col = ACCENT_MAP.get(ac, ORANGE)
        tb  = slide.shapes.add_textbox(int(left), int(cursor_y), int(width), int(Inches(0.9)))
        tf  = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = value_el.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(calc_kpi_value_pt(width)); run.font.color.rgb = col
        cursor_y += Inches(0.92)

    label_el = card_el.select_one('.kpi-label')
    if label_el:
        tb = slide.shapes.add_textbox(int(left+Inches(0.1)), int(cursor_y),
                                       int(width-Inches(0.2)), int(Inches(0.32)))
        tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = label_el.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(12); run.font.color.rgb = DARK
        cursor_y += Inches(0.34)

    desc_el = card_el.select_one('.kpi-desc')
    if desc_el:
        tb = slide.shapes.add_textbox(int(left+Inches(0.1)), int(cursor_y),
                                       int(width-Inches(0.2)), int(Inches(0.3)))
        tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = desc_el.get_text(strip=True)
        run.font.name = 'Calibri'; run.font.size = Pt(11); run.font.color.rgb = MGRAY

def render_vs_box(slide, card_el, left, top, width, height, accent='orange'):
    accent_color = ACCENT_MAP.get(accent, ORANGE)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(width), int(height))
    card.fill.solid(); card.fill.fore_color.rgb = BGRAY
    card.line.color.rgb = accent_color; card.line.width = Pt(2)

    cursor_y = top + Inches(0.18)

    h3 = card_el.select_one('h3')
    if h3:
        tb = slide.shapes.add_textbox(
            int(left + Inches(0.18)), int(cursor_y),
            int(width - Inches(0.36)), int(Inches(0.42)))
        tf = tb.text_frame; p = tf.paragraphs[0]
        run = p.add_run(); run.text = h3.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(calc_heading_pt(width)); run.font.color.rgb = accent_color
        cursor_y += Inches(0.5)

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(left + Inches(0.18)), int(cursor_y - Inches(0.08)),
            int(width - Inches(0.36)), int(Pt(1.5)))
        line.fill.solid(); line.fill.fore_color.rgb = LGRAY; line.line.fill.background()

    bullets = card_el.select('ul li')
    if bullets:
        avail_h   = top + height - cursor_y - Inches(0.15)
        bullet_pt = calc_bullet_pt(avail_h, len(bullets))
        space_pt  = max(1.0, bullet_pt * 0.18)
        tb = slide.shapes.add_textbox(
            int(left + Inches(0.2)), int(cursor_y),
            int(width - Inches(0.35)), int(avail_h))
        tf = tb.text_frame; tf.word_wrap = True
        for j, li in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run(); run.text = '\u2022  ' + li.get_text(strip=True)
            run.font.name = 'Calibri'; run.font.size = Pt(bullet_pt)
            run.font.color.rgb = DARK; p.space_after = Pt(space_pt)

def render_roadmap(slide, content_el):
    phase_row = content_el.select_one('.phase-row')
    phases    = phase_row.select('.phase') if phase_row else []
    n = len(phases)
    if n == 0: return

    LEFT    = Inches(0.5); TOP = Inches(1.4)
    TOTAL_W = Inches(12.3); ARROW_H = Inches(0.95); GAP = Inches(0.07)
    arrow_w = (TOTAL_W - GAP * (n - 1)) / n
    COLOR_MAP = {'orange': ORANGE, 'gold': GOLD, 'dark': DARK}

    for i, phase in enumerate(phases):
        px         = int(LEFT + i * (arrow_w + GAP))
        color_key  = phase.get('data-color', 'orange')
        fill_color = COLOR_MAP.get(color_key, ORANGE)
        text_color = DARK if color_key == 'gold' else WHITE

        arrow = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, px, int(TOP), int(arrow_w), int(ARROW_H))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = fill_color
        arrow.line.fill.background()

        tf = arrow.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = phase.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(calc_heading_pt(arrow_w, max_pt=13.0, min_pt=10.0))
        run.font.color.rgb = text_color

    details_row  = content_el.select_one('.phase-details')
    if not details_row: return
    detail_items = details_row.select('.phase-detail')
    detail_top   = TOP + ARROW_H + Inches(0.22)
    detail_h     = Inches(7.5) - detail_top - Inches(0.45)

    for i, detail in enumerate(detail_items[:n]):
        dx = int(LEFT + i * (arrow_w + GAP))
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, dx, int(detail_top), int(arrow_w), int(detail_h))
        bg.fill.solid(); bg.fill.fore_color.rgb = BGRAY
        bg.line.color.rgb = LGRAY; bg.line.width = Pt(1)

        bullets = detail.select('li')
        if bullets:
            avail_detail = detail_h - Inches(0.24)
            detail_pt    = calc_bullet_pt(avail_detail, len(bullets), max_pt=13.5, min_pt=9.5)
            space_pt     = max(1.0, detail_pt * 0.2)
            tb = slide.shapes.add_textbox(
                int(dx + Inches(0.15)), int(detail_top + Inches(0.12)),
                int(arrow_w - Inches(0.3)), int(avail_detail))
            tf = tb.text_frame; tf.word_wrap = True
            for j, li in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run = p.add_run(); run.text = '\u2022  ' + li.get_text(strip=True)
                run.font.name = 'Calibri'; run.font.size = Pt(detail_pt)
                run.font.color.rgb = DARK; p.space_after = Pt(space_pt)

def render_process_flow(slide, content_el):
    steps = content_el.select('.step')
    n = len(steps)
    if n == 0: return

    LEFT    = Inches(0.5); TOP = Inches(1.4)
    TOTAL_W = Inches(12.3); STEP_H = Inches(1.15); ARROW_W = Inches(0.28)
    step_w  = (TOTAL_W - ARROW_W * (n - 1)) / n

    for i, step in enumerate(steps):
        sx = LEFT + i * (step_w + ARROW_W)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(sx), int(TOP), int(step_w), int(STEP_H))
        box.fill.solid(); box.fill.fore_color.rgb = ORANGE; box.line.fill.background()

        num_text = step.get('data-num', str(i + 1))
        tb = slide.shapes.add_textbox(
            int(sx + Inches(0.1)), int(TOP + Inches(0.05)),
            int(Inches(0.42)), int(STEP_H - Inches(0.1)))
        tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p  = tf.paragraphs[0]
        run = p.add_run(); run.text = str(num_text)
        run.font.name = 'Poppins'; run.font.bold = True
        run.font.size = Pt(24); run.font.color.rgb = WHITE

        tx   = sx + Inches(0.56); tw = step_w - Inches(0.66)
        h3   = step.select_one('h3')
        p_el = step.select_one('p')

        if h3:
            tb2 = slide.shapes.add_textbox(
                int(tx), int(TOP + Inches(0.1)), int(tw), int(Inches(0.44)))
            tf2 = tb2.text_frame
            p2  = tf2.paragraphs[0]
            run2 = p2.add_run(); run2.text = h3.get_text(strip=True)
            run2.font.name = 'Poppins'; run2.font.bold = True
            run2.font.size = Pt(13); run2.font.color.rgb = WHITE

        if p_el:
            tb3 = slide.shapes.add_textbox(
                int(tx), int(TOP + Inches(0.54)), int(tw), int(Inches(0.52)))
            tf3 = tb3.text_frame; tf3.word_wrap = True
            p3  = tf3.paragraphs[0]
            run3 = p3.add_run(); run3.text = p_el.get_text(strip=True)
            run3.font.name = 'Calibri'; run3.font.size = Pt(11)
            run3.font.color.rgb = RGBColor(0xFF, 0xE0, 0xCC)

        if i < n - 1:
            ax = sx + step_w + Inches(0.04)
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, int(ax), int(TOP + STEP_H/2 - Inches(0.12)),
                int(ARROW_W - Inches(0.08)), int(Inches(0.24)))
            arr.fill.solid(); arr.fill.fore_color.rgb = DARK; arr.line.fill.background()

    detail_row = content_el.select_one('.step-details-row')
    if not detail_row: return
    details    = detail_row.select('.step-detail')
    dt_top     = TOP + STEP_H + Inches(0.22)
    dt_h       = Inches(7.5) - dt_top - Inches(0.45)

    for i, detail in enumerate(details[:n]):
        dx = LEFT + i * (step_w + ARROW_W)
        tb = slide.shapes.add_textbox(int(dx), int(dt_top), int(step_w), int(dt_h))
        tf = tb.text_frame; tf.word_wrap = True
        for j, li in enumerate(detail.select('li')):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run(); run.text = '\u2022  ' + li.get_text(strip=True)
            run.font.name = 'Calibri'; run.font.size = Pt(11)
            run.font.color.rgb = DARK; p.space_after = Pt(3)

def render_next_steps(slide, content_el):
    steps  = content_el.select('.next-step')
    n      = len(steps)
    if n == 0: return

    LEFT    = Inches(0.5); TOP = Inches(1.42)
    TOTAL_H = Inches(5.5); GAP = Inches(0.14)
    step_h  = (TOTAL_H - GAP * (n - 1)) / n

    for i, step in enumerate(steps):
        sy           = TOP + i * (step_h + GAP)
        accent       = step.get('data-accent', 'orange')
        accent_color = ACCENT_MAP.get(accent, ORANGE)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(Inches(0.5)), int(sy), int(Inches(12.3)), int(step_h))
        bg.fill.solid(); bg.fill.fore_color.rgb = BGRAY; bg.line.fill.background()

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(Inches(0.5)), int(sy), int(Pt(5)), int(step_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent_color; bar.line.fill.background()

        date_el = step.select_one('.next-step-date')
        if date_el:
            tb = slide.shapes.add_textbox(
                int(Inches(0.75)), int(sy + step_h/2 - Inches(0.18)),
                int(Inches(1.1)), int(Inches(0.36)))
            tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p  = tf.paragraphs[0]
            run = p.add_run(); run.text = date_el.get_text(strip=True)
            run.font.name = 'Poppins'; run.font.bold = True
            run.font.size = Pt(13); run.font.color.rgb = accent_color

        content_el2 = step.select_one('.next-step-content')
        if content_el2:
            h3   = content_el2.select_one('h3')
            p_el = content_el2.select_one('p')
            cx   = Inches(2.0); cw = Inches(10.6)
            cy   = sy + step_h/2 - (Inches(0.2) if p_el else Inches(0.12))

            if h3:
                tb = slide.shapes.add_textbox(int(cx), int(cy), int(cw), int(Inches(0.32)))
                tf = tb.text_frame
                p  = tf.paragraphs[0]
                run = p.add_run(); run.text = h3.get_text(strip=True)
                run.font.name = 'Poppins'; run.font.bold = True
                run.font.size = Pt(13); run.font.color.rgb = DARK
                cy += Inches(0.32)

            if p_el:
                tb = slide.shapes.add_textbox(int(cx), int(cy), int(cw), int(Inches(0.26)))
                tf = tb.text_frame; tf.word_wrap = True
                p  = tf.paragraphs[0]
                run = p.add_run(); run.text = p_el.get_text(strip=True)
                run.font.name = 'Calibri'; run.font.size = Pt(11); run.font.color.rgb = MGRAY

def render_cover(slide, slide_div):
    cc = slide_div.select_one('.cover-center')
    if not cc: return

    title_el    = cc.select_one('.cover-title')
    subtitle_el = cc.select_one('.cover-subtitle')

    if title_el:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.0))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for child in title_el.children:
            if isinstance(child, NavigableString):
                t = str(child)
                if not t.strip(): continue
                run = p.add_run(); run.text = t
                run.font.name = 'Poppins'; run.font.bold = True
                run.font.size = Pt(40); run.font.color.rgb = DARK
            elif isinstance(child, Tag):
                run = p.add_run(); run.text = child.get_text()
                run.font.name = 'Poppins'; run.font.bold = True; run.font.size = Pt(40)
                run.font.color.rgb = ORANGE if 'accent' in child.get('class', []) else DARK

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.17), Inches(4.32), Inches(3.0), Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()

    if subtitle_el:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.7))
        tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = subtitle_el.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.size = Pt(18); run.font.color.rgb = MGRAY

def render_divider(slide, slide_div):
    dc = slide_div.select_one('.divider-center')
    if not dc: return
    title_el    = dc.select_one('.divider-title')
    subtitle_el = dc.select_one('.divider-subtitle')

    if title_el:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.8))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for child in title_el.children:
            if isinstance(child, NavigableString):
                t = str(child)
                if not t.strip(): continue
                run = p.add_run(); run.text = t
                run.font.name = 'Poppins'; run.font.bold = True
                run.font.size = Pt(36); run.font.color.rgb = WHITE
            elif isinstance(child, Tag):
                run = p.add_run(); run.text = child.get_text()
                run.font.name = 'Poppins'; run.font.bold = True; run.font.size = Pt(36)
                run.font.color.rgb = ORANGE if 'accent' in child.get('class', []) else WHITE

    if subtitle_el:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.55))
        tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = subtitle_el.get_text(strip=True)
        run.font.name = 'Poppins'; run.font.size = Pt(18); run.font.color.rgb = LGRAY

def render_content(slide, content_el):
    classes = content_el.get('class', [])

    if 'roadmap' in classes or content_el.select_one('.phase-row'):
        render_roadmap(slide, content_el)
        return

    if 'process-flow' in classes or content_el.select_one('.process-steps'):
        render_process_flow(slide, content_el)
        return

    if 'next-steps' in classes or content_el.select_one('.next-step'):
        render_next_steps(slide, content_el)
        return

    n_cols = 1
    for c in classes:
        if c.startswith('columns-'):
            try: n_cols = int(c.split('-')[1])
            except: pass

    LEFT    = Inches(0.5); TOP = Inches(1.4)
    TOTAL_W = Inches(12.3); TOTAL_H = Inches(5.5); GAP = Inches(0.18)
    col_w   = (TOTAL_W - GAP * (n_cols - 1)) / n_cols

    items = content_el.select(':scope > .card, :scope > .kpi-card, :scope > .vs-box')

    for i, item in enumerate(items):
        if i >= n_cols: break
        cx = LEFT + i * (col_w + GAP)
        ic = item.get('class', [])
        if 'kpi-card' in ic:
            render_kpi_card(slide, item, cx, TOP, col_w, TOTAL_H)
        elif 'vs-box' in ic:
            accent = item.get('data-accent', 'orange')
            render_vs_box(slide, item, cx, TOP, col_w, TOTAL_H, accent)
        else:
            accent = item.get('data-accent', 'orange')
            render_card(slide, item, cx, TOP, col_w, TOTAL_H, accent)

# === HAUPTFUNKTION ===
def html_to_pptx(html_content: str, output_path: str) -> str:
    soup     = BeautifulSoup(html_content, 'html.parser')
    prs      = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page_num = 0

    for slide_div in soup.select('.slide'):
        page_num += 1
        slide   = prs.slides.add_slide(prs.slide_layouts[6])
        classes = slide_div.get('class', [])
        is_cover   = 'slide-cover'   in classes
        is_divider = 'slide-divider' in classes

        if is_divider:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()

        add_logo(slide, dark_bg=is_divider)
        add_gradient_line(slide)

        if is_cover:
            render_cover(slide, slide_div)
            add_footer(slide, page_num)
            continue

        if is_divider:
            render_divider(slide, slide_div)
            add_footer(slide, page_num)
            continue

        title_el = slide_div.select_one('.slide-title')
        if title_el:
            add_slide_title(slide, title_el)

        content_el = slide_div.select_one('.slide-content')
        if content_el:
            render_content(slide, content_el)

        add_footer(slide, page_num)

    prs.save(output_path)
    return output_path

# === AUSFUEHRUNGS-LOGIK ===
import re

# Logo-Pfade als Base64 in HTML ersetzen (damit HTML-Datei standalone ist)
with open(LOGO_PATH, 'rb')  as f: logo_b64       = base64.b64encode(f.read()).decode()
with open(LOGO_WHITE, 'rb') as f: logo_white_b64 = base64.b64encode(f.read()).decode()

HTML_CONTENT = HTML_CONTENT.replace(
    'LOGO_SRC',       f'data:image/png;base64,{logo_b64}')
HTML_CONTENT = HTML_CONTENT.replace(
    'LOGO_WHITE_SRC', f'data:image/png;base64,{logo_white_b64}')

# Dateinamen ableiten (Slug aus Thema)
slug = re.sub(r'[^a-z0-9]+', '_', TOPIC.lower()).strip('_')[:40]

# HTML speichern
html_path = Path.cwd() / f"presentation_{slug}.html"
html_path.write_text(HTML_CONTENT, encoding='utf-8')
print(f"HTML-Preview: {html_path}")

# Im Browser oeffnen
subprocess.Popen(['open', str(html_path)])

# PPTX erzeugen
pptx_path = str(Path.cwd() / f"presentation_{slug}.pptx")
html_to_pptx(HTML_CONTENT, pptx_path)
print(f"PowerPoint:   {pptx_path}")

# PPTX oeffnen
subprocess.Popen(['open', pptx_path])
